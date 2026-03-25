"""
Production quality gate for editable PDF -> PPTX conversion.

Usage:
  python quality_gate_pdf_to_pptx.py --pdf input.pdf --pptx output.pptx
  python quality_gate_pdf_to_pptx.py --pdf input.pdf --convert --out output.pptx
"""

import argparse
import json
import os
from typing import Dict

import fitz
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from converter import pdf_to_editable_pptx


def _safe_ratio(n: int, d: int) -> float:
    if d <= 0:
        return 1.0
    return float(n) / float(d)


def analyze_pdf(pdf_path: str) -> Dict[str, int]:
    out = {
        "pages": 0,
        "text_spans": 0,
        "image_instances": 0,
        "vector_drawings": 0,
        "table_regions": 0,
    }
    with fitz.open(pdf_path) as doc:
        out["pages"] = len(doc)
        for page in doc:
            text = page.get_text("dict")
            for block in text.get("blocks", []):
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        t = span.get("text", "")
                        if t:
                            out["text_spans"] += 1

            for img in page.get_images(full=True):
                out["image_instances"] += len(page.get_image_rects(img))

            for dr in page.get_drawings():
                if dr.get("fill") or dr.get("color"):
                    out["vector_drawings"] += 1

            try:
                finder = page.find_tables()
                out["table_regions"] += len(finder.tables)
            except Exception:
                pass
    return out


def analyze_pptx(pptx_path: str) -> Dict[str, int]:
    out = {
        "slides": 0,
        "text_shapes": 0,
        "text_runs": 0,
        "pictures": 0,
        "vectors": 0,
        "tables": 0,
    }
    prs = Presentation(pptx_path)
    out["slides"] = len(prs.slides)

    for slide in prs.slides:
        for shape in slide.shapes:
            st = shape.shape_type
            if st == MSO_SHAPE_TYPE.PICTURE:
                out["pictures"] += 1
            elif st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.LINE):
                out["vectors"] += 1
            elif st == MSO_SHAPE_TYPE.TABLE:
                out["tables"] += 1

            if hasattr(shape, "text_frame") and shape.has_text_frame:
                text = shape.text_frame.text or ""
                if text.strip():
                    out["text_shapes"] += 1
                for para in shape.text_frame.paragraphs:
                    out["text_runs"] += len(para.runs)
    return out


def build_report(pdf_stats: Dict[str, int], pptx_stats: Dict[str, int]) -> Dict[str, object]:
    coverage = {
        "page_ratio": _safe_ratio(pptx_stats["slides"], pdf_stats["pages"]),
        "text_run_ratio": _safe_ratio(pptx_stats["text_runs"], pdf_stats["text_spans"]),
        "image_ratio": _safe_ratio(pptx_stats["pictures"], pdf_stats["image_instances"]),
        "table_ratio": _safe_ratio(pptx_stats["tables"], pdf_stats["table_regions"]),
        "vector_ratio": _safe_ratio(pptx_stats["vectors"], pdf_stats["vector_drawings"]),
    }

    # Table detection is inherently noisier than text/image extraction.
    table_threshold = 0.70 if pdf_stats["table_regions"] <= 3 else 0.30

    checks = {
        "pages_ok": coverage["page_ratio"] >= 1.0,
        "text_ok": coverage["text_run_ratio"] >= 0.90,
        "images_ok": coverage["image_ratio"] >= 0.90,
        "tables_ok": coverage["table_ratio"] >= table_threshold,
        "vectors_ok": coverage["vector_ratio"] >= 0.50,
    }

    production_ready = all(checks.values())
    return {
        "production_ready": production_ready,
        "checks": checks,
        "coverage": coverage,
        "pdf_stats": pdf_stats,
        "pptx_stats": pptx_stats,
    }


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--pdf", required=True, help="Input PDF path")
    ap.add_argument("--pptx", help="Existing output PPTX path")
    ap.add_argument("--convert", action="store_true", help="Run conversion before scoring")
    ap.add_argument("--out", help="Output PPTX path when using --convert")
    ap.add_argument("--report", help="Optional JSON report output path")
    args = ap.parse_args()

    pdf_path = os.path.abspath(args.pdf)
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF not found: {pdf_path}")

    pptx_path = args.pptx
    if args.convert:
        pptx_path = os.path.abspath(args.out or (os.path.splitext(pdf_path)[0] + ".editable.pptx"))
        pdf_to_editable_pptx(pdf_path, pptx_path, on_status=lambda m: print(f"[convert] {m}"))
    elif not pptx_path:
        raise ValueError("Provide --pptx for scoring, or use --convert with optional --out")

    pptx_path = os.path.abspath(pptx_path)
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    pdf_stats = analyze_pdf(pdf_path)
    pptx_stats = analyze_pptx(pptx_path)
    report = build_report(pdf_stats, pptx_stats)

    print(json.dumps(report, indent=2))
    if args.report:
        with open(args.report, "w", encoding="utf-8") as f:
            json.dump(report, f, indent=2)

    if not report["production_ready"]:
        raise SystemExit(2)


if __name__ == "__main__":
    main()
