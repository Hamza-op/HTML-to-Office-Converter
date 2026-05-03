import os
from io import BytesIO
from typing import List, Callable, Optional
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from scene_graph import Page, Image, TextBlock, Table, VectorPath, UnknownRegion, Group

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY
}

try:
    from pptx.enum.shapes import MSO_SHAPE
    RECTANGLE = MSO_SHAPE.RECTANGLE
except ImportError:
    RECTANGLE = 1

def _set_fill(shape, rgb):
    if rgb is None:
        shape.fill.background()
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)

def _set_line(shape, rgb, width_pt):
    if rgb is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = RGBColor(*rgb)
    if width_pt and width_pt > 0:
        shape.line.width = Pt(max(width_pt, 0.25))

def pts_close(pt1, pt2):
    if pt1 is None or pt2 is None: return False
    return abs(pt1[0] - pt2[0]) < 1e-3 and abs(pt1[1] - pt2[1]) < 1e-3

def write_scene_graph_to_pptx(
    pages: List[Page],
    output_path: str,
    slide_size: str = "Widescreen (16:9)",
    warnings: Optional[List[str]] = None,
    on_status: Optional[Callable[[str], None]] = None
) -> str:
    status = on_status or (lambda _: None)
    status("Writing scene graph to editable PPTX...")
    if warnings is None:
        warnings = []

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    metrics = {
        "pages": len(pages),
        "text_boxes": 0,
        "text_runs": 0,
        "images": 0,
        "vectors": 0,
        "tables": 0
    }

    if pages:
        prs.slide_width = Pt(max(pages[0].width, 10))
        prs.slide_height = Pt(max(pages[0].height, 10))
    else:
        if "4:3" in slide_size:
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
        else:
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

    for page in pages:
        status(f"Writing page {page.page_num}/{len(pages)}...")
        slide = prs.slides.add_slide(blank_layout)
        
        bg_shape = slide.shapes.add_shape(
            RECTANGLE, Pt(0), Pt(0), Pt(page.width), Pt(page.height)
        )
        _set_fill(bg_shape, (255, 255, 255))
        bg_shape.line.fill.background()
        
        # We process children of the root group
        emitted_count = 0
        
        def iter_nodes(parent):
            for n in sorted(parent.children, key=lambda x: getattr(x, 'z_order', 0)):
                if isinstance(n, Group):
                    yield from iter_nodes(n)
                else:
                    yield n
                    
        for node in iter_nodes(page.root):
            if getattr(node, "confidence", 1.0) < 0.5:
                warnings.append(f"Page {page.page_num}: Low confidence ({node.confidence}) on {getattr(node, 'provenance', 'object')}")
                
            if isinstance(node, VectorPath):
                # Check if it's a full page background
                w_cov = (node.bbox[2] - node.bbox[0]) / max(page.width, 1) if node.bbox else 0
                h_cov = (node.bbox[3] - node.bbox[1]) / max(page.height, 1) if node.bbox else 0
                
                # If only one segment "re" and covers full page
                if len(node.segments) == 1 and node.segments[0]["type"] == "re":
                    if w_cov >= 0.98 and h_cov >= 0.98 and node.fill_color:
                        _set_fill(bg_shape, node.fill_color)
                        continue
                
                builder = None
                last_pt = None
                
                for seg in node.segments:
                    op = seg["type"]
                    args = seg["args"]
                    if op == "l":
                        p0, p1 = args[0], args[1]
                        if builder is None:
                            builder = slide.shapes.build_freeform(Pt(p0.x), Pt(p0.y))
                        elif not pts_close(last_pt, (p0.x, p0.y)):
                            builder.move_to(Pt(p0.x), Pt(p0.y))
                        builder.add_line_segments([(Pt(p1.x), Pt(p1.y))], close=False)
                        last_pt = (p1.x, p1.y)
                    elif op == "c":
                        p0, p1, p2, p3 = args[0], args[1], args[2], args[3]
                        if builder is None:
                            builder = slide.shapes.build_freeform(Pt(p0.x), Pt(p0.y))
                        elif not pts_close(last_pt, (p0.x, p0.y)):
                            builder.move_to(Pt(p0.x), Pt(p0.y))
                        pts = []
                        for step in range(1, 11):
                            t = step / 10.0
                            mt = 1 - t
                            x = (mt**3)*p0.x + 3*(mt**2)*t*p1.x + 3*mt*(t**2)*p2.x + (t**3)*p3.x
                            y = (mt**3)*p0.y + 3*(mt**2)*t*p1.y + 3*mt*(t**2)*p2.y + (t**3)*p3.y
                            pts.append((Pt(x), Pt(y)))
                        builder.add_line_segments(pts, close=False)
                        last_pt = (p3.x, p3.y)
                    elif op == "re":
                        rect = args[0]
                        if builder is None:
                            builder = slide.shapes.build_freeform(Pt(rect.x0), Pt(rect.y0))
                        else:
                            builder.move_to(Pt(rect.x0), Pt(rect.y0))
                        builder.add_line_segments([
                            (Pt(rect.x1), Pt(rect.y0)),
                            (Pt(rect.x1), Pt(rect.y1)),
                            (Pt(rect.x0), Pt(rect.y1)),
                            (Pt(rect.x0), Pt(rect.y0))
                        ], close=True)
                        last_pt = (rect.x0, rect.y0)
                    elif op == "qu":
                        q = args[0]
                        if builder is None:
                            builder = slide.shapes.build_freeform(Pt(q.ul.x), Pt(q.ul.y))
                        else:
                            builder.move_to(Pt(q.ul.x), Pt(q.ul.y))
                        builder.add_line_segments([
                            (Pt(q.ur.x), Pt(q.ur.y)),
                            (Pt(q.lr.x), Pt(q.lr.y)),
                            (Pt(q.ll.x), Pt(q.ll.y)),
                            (Pt(q.ul.x), Pt(q.ul.y))
                        ], close=True)
                        last_pt = (q.ul.x, q.ul.y)
                
                if builder is not None:
                    try:
                        shape = builder.convert_to_shape()
                        _set_fill(shape, node.fill_color)
                        _set_line(shape, node.stroke_color, node.stroke_width)
                        metrics["vectors"] += 1
                        emitted_count += 1
                    except Exception as e:
                        warnings.append(f"Page {page.page_num}: Failed to write {getattr(node, 'provenance', 'unknown')} vector path ({e})")
                        metrics.setdefault("failed_vectors", 0)
                        metrics["failed_vectors"] += 1
                        
            elif isinstance(node, Image):
                try:
                    slide.shapes.add_picture(
                        BytesIO(node.image_bytes),
                        Pt(max(node.bbox[0], 0)),
                        Pt(max(node.bbox[1], 0)),
                        width=Pt(max(node.bbox[2]-node.bbox[0], 4)),
                        height=Pt(max(node.bbox[3]-node.bbox[1], 4))
                    )
                    metrics["images"] += 1
                    emitted_count += 1
                except Exception as e:
                    warnings.append(f"Page {page.page_num}: Failed to write {getattr(node, 'provenance', 'unknown')} image ({e})")
                    metrics.setdefault("failed_images", 0)
                    metrics["failed_images"] += 1
                    
            elif isinstance(node, TextBlock):
                left, top, right, bottom = node.bbox
                try:
                    text_box = slide.shapes.add_textbox(
                        Pt(max(left, 0)), Pt(max(top, 0)), 
                        Pt(max(right-left, 4)), Pt(max(bottom-top, 4))
                    )
                    text_frame = text_box.text_frame
                    text_frame.clear()
                    text_frame.word_wrap = True
                    text_frame.margin_bottom = 0
                    text_frame.margin_top = 0
                    text_frame.margin_left = 0
                    text_frame.margin_right = 0
                    
                    first_para = True
                    for para_data in node.paragraphs:
                        if first_para:
                            para = text_frame.paragraphs[0]
                            first_para = False
                        else:
                            para = text_frame.add_paragraph()
                        
                        if para_data.alignment in ALIGN_MAP:
                            para.alignment = ALIGN_MAP[para_data.alignment]
                        
                        for run_data in para_data.runs:
                            run = para.add_run()
                            run.text = run_data.text
                            font = run.font
                            if run_data.font_name: font.name = run_data.font_name
                            if run_data.font_size: font.size = Pt(max(float(run_data.font_size), 1.0))
                            if run_data.color: font.color.rgb = RGBColor(*run_data.color)
                            font.bold = run_data.bold
                            font.italic = run_data.italic
                            metrics["text_runs"] += 1
                    metrics["text_boxes"] += 1
                    emitted_count += 1
                except Exception as e:
                    warnings.append(f"Page {page.page_num}: Failed to write {getattr(node, 'provenance', 'unknown')} text block ({e})")
                    metrics.setdefault("failed_text_boxes", 0)
                    metrics["failed_text_boxes"] += 1
                    
            elif isinstance(node, UnknownRegion) and node.image_fallback:
                try:
                    slide.shapes.add_picture(
                        BytesIO(node.image_fallback),
                        Pt(max(node.bbox[0], 0)),
                        Pt(max(node.bbox[1], 0)),
                        width=Pt(max(node.bbox[2]-node.bbox[0], 4)),
                        height=Pt(max(node.bbox[3]-node.bbox[1], 4))
                    )
                    metrics["images"] += 1
                    metrics.setdefault("raster_fallback_pages", 0)
                    metrics["raster_fallback_pages"] += 1
                    emitted_count += 1
                except Exception as e:
                    warnings.append(f"Page {page.page_num}: Failed to write {getattr(node, 'provenance', 'unknown')} raster fallback ({e})")
                    
            elif isinstance(node, Table):
                try:
                    left, top, right, bottom = node.bbox
                    shape = slide.shapes.add_table(
                        node.rows, node.cols,
                        Pt(left), Pt(top), Pt(right-left), Pt(bottom-top)
                    )
                    table = shape.table
                    for cell in node.cells:
                        if cell.row < node.rows and cell.col < node.cols:
                            tcell = table.cell(cell.row, cell.col)
                            if cell.text_block and cell.text_block.paragraphs:
                                tcell.text = "".join(r.text for p in cell.text_block.paragraphs for r in p.runs)
                                para = tcell.text_frame.paragraphs[0]
                                para.font.size = Pt(10)
                                if any(r.bold for p in cell.text_block.paragraphs for r in p.runs):
                                    para.font.bold = True
                    metrics["tables"] += 1
                    emitted_count += 1
                except Exception as e:
                    warnings.append(f"Page {page.page_num}: Failed to write {getattr(node, 'provenance', 'unknown')} table ({e})")
                    metrics.setdefault("failed_tables", 0)
                    metrics["failed_tables"] += 1
                    
        if emitted_count == 0 and len(page.root.children) > 0:
            warnings.append(f"Page {page.page_num}: Wrote 0 objects, page is effectively blank! Applying fallback.")
            if getattr(page, 'fallback_bytes', None):
                try:
                    slide.shapes.add_picture(
                        BytesIO(page.fallback_bytes),
                        Pt(0), Pt(0),
                        width=Pt(page.width), height=Pt(page.height)
                    )
                    metrics["images"] += 1
                    emitted_count += 1
                    metrics.setdefault("raster_fallback_pages", 0)
                    metrics["raster_fallback_pages"] += 1
                except Exception as e:
                    warnings.append(f"Page {page.page_num}: Final fallback failed ({e})")

    prs.save(output_path)
    fallback_pages = metrics.get("raster_fallback_pages", 0)
    status(
        "Editable PPTX metrics:\n"
        f" - Pages: {metrics['pages']}\n"
        f" - Text Boxes: {metrics['text_boxes']} emitted, {metrics.get('failed_text_boxes', 0)} failed\n"
        f" - Text Runs: {metrics['text_runs']}\n"
        f" - Images: {metrics['images']} emitted, {metrics.get('failed_images', 0)} failed\n"
        f" - Vectors: {metrics['vectors']} emitted, {metrics.get('failed_vectors', 0)} failed\n"
        f" - Tables: {metrics['tables']} emitted, {metrics.get('failed_tables', 0)} failed\n"
        f" - Fallback Pages: {fallback_pages}"
    )
    if warnings:
        status(f"Warnings: {len(warnings)}")
        for w in warnings:
            status(f" - {w}")
    status(f"Editable PPTX saved: {os.path.basename(output_path)}")
    return output_path
