"""
Core conversion engine.
HTML → Playwright (render) → PDF → Editable DOCX (via pdf2docx)
HTML → Playwright (render) → Screenshots → PPTX
"""

import asyncio
import math
import os
import sys
import platform
import tempfile
from pathlib import Path
from typing import Callable, Optional

from PIL import Image
from io import BytesIO

# ── Page size presets (width_inches, height_inches) ──
PAGE_SIZES = {
    "A4": (8.27, 11.69),
    "Letter": (8.5, 11.0),
    "Legal": (8.5, 14.0),
    "A3": (11.69, 16.54),
    "A5": (5.83, 8.27),
}

# ── Orientation ──
ORIENTATIONS = ["Portrait", "Landscape"]

# ── DPI presets ──
DPI_PRESETS = {
    "Standard (150 DPI)": 150,
    "High (200 DPI)": 200,
    "Ultra (300 DPI)": 300,
}


def _build_file_url(abs_path: str) -> str:
    """Build a file:// URL that works across platforms."""
    if platform.system() == "Windows":
        return "file:///" + abs_path.replace("\\", "/")
    return "file://" + abs_path


# Channels to try in order: system Chrome, Edge, Firefox, then Playwright Chromium
_BROWSER_CHANNELS = ["chrome", "msedge", None]


async def _launch_best_browser(pw, **kwargs):
    """Try system Chrome → Edge → Playwright Chromium in order."""
    last_err = None
    for channel in _BROWSER_CHANNELS:
        try:
            opts = dict(kwargs)
            if channel:
                opts["channel"] = channel
            return await pw.chromium.launch(**opts)
        except Exception as e:
            last_err = e
            continue
    raise RuntimeError(
        f"No usable browser found. Install Chrome or Edge, or run: "
        f"playwright install chromium\n({last_err})"
    )


# ─────────────────────────────────────────────────────────────────
#  HTML → PDF (via Playwright)
# ─────────────────────────────────────────────────────────────────

async def _render_html_to_pdf(
    html_path: str,
    pdf_path: str,
    page_size: str = "A4",
    orientation: str = "Portrait",
    margin_inches: float = 0.5,
    on_status: Optional[Callable[[str], None]] = None,
) -> str:
    """
    Render HTML in a headless browser and export as PDF.
    The PDF preserves all styling, fonts, and layout perfectly.
    """
    from playwright.async_api import async_playwright

    status = on_status or (lambda _: None)

    pw_inches, ph_inches = PAGE_SIZES.get(page_size, PAGE_SIZES["A4"])
    if orientation == "Landscape":
        pw_inches, ph_inches = ph_inches, pw_inches

    status("Launching browser...")
    async with async_playwright() as pw:
        browser = await _launch_best_browser(pw, headless=True)
        page = await browser.new_page()

        abs_path = os.path.abspath(html_path)
        file_url = _build_file_url(abs_path)

        status("Loading HTML...")
        await page.goto(file_url, wait_until="networkidle")
        await page.evaluate("() => document.fonts.ready")

        # Force screen media so @media print rules don't strip backgrounds
        await page.emulate_media(media="screen")

        # Remove @media print rules and floating UI elements
        await page.evaluate("""() => {
            // Strip all @media print CSS rules that hide backgrounds
            for (const sheet of document.styleSheets) {
                try {
                    const rules = sheet.cssRules || [];
                    for (let i = rules.length - 1; i >= 0; i--) {
                        if (rules[i].type === CSSRule.MEDIA_RULE &&
                            rules[i].conditionText === 'print') {
                            sheet.deleteRule(i);
                        }
                    }
                } catch(e) {}  // skip cross-origin sheets
            }

            // Remove no-print elements (floating buttons etc.)
            const selectors = ['.fab-container', '.no-print', '[data-no-export]'];
            selectors.forEach(sel => {
                document.querySelectorAll(sel).forEach(el => el.remove());
            });
        }""")

        margin_str = f"{margin_inches}in"

        status("Generating PDF...")
        await page.pdf(
            path=pdf_path,
            width=f"{pw_inches}in",
            height=f"{ph_inches}in",
            margin={
                "top": margin_str,
                "bottom": margin_str,
                "left": margin_str,
                "right": margin_str,
            },
            print_background=True,
            prefer_css_page_size=True,
        )

        await browser.close()

    status("PDF generated.")
    return pdf_path


def render_to_pdf(
    html_path: str,
    pdf_path: str,
    page_size: str = "A4",
    orientation: str = "Portrait",
    margin_inches: float = 0.5,
    on_status: Optional[Callable[[str], None]] = None,
) -> str:
    """Synchronous wrapper for PDF rendering."""
    return asyncio.run(
        _render_html_to_pdf(
            html_path,
            pdf_path,
            page_size,
            orientation,
            margin_inches,
            on_status,
        )
    )


# ─────────────────────────────────────────────────────────────────
#  PDF → Editable DOCX (via pdf2docx)
# ─────────────────────────────────────────────────────────────────

def pdf_to_docx(
    pdf_path: str,
    docx_path: str,
    on_status: Optional[Callable[[str], None]] = None,
) -> str:
    """
    Convert PDF to editable DOCX using pdf2docx.
    Preserves text, tables, images, fonts, and layout.
    """
    from pdf2docx import Converter

    status = on_status or (lambda _: None)

    status("Converting PDF to editable DOCX...")
    cv = Converter(pdf_path)
    cv.convert(docx_path)
    cv.close()

    status(f"DOCX saved: {os.path.basename(docx_path)}")
    return docx_path


def pdf_to_images(
    pdf_path: str,
    dpi: int = 150,
    on_status: Optional[Callable[[str], None]] = None,
) -> list:
    """
    Render PDF pages as PIL Images for the preview viewer.
    """
    import fitz  # PyMuPDF

    status = on_status or (lambda _: None)
    status("Rendering preview pages...")

    doc = fitz.open(pdf_path)
    images = []

    for i in range(len(doc)):
        page = doc[i]
        pix = page.get_pixmap(dpi=dpi)
        img = Image.open(BytesIO(pix.tobytes("png")))
        images.append(img)

    doc.close()
    status(f"Rendered {len(images)} preview page(s).")
    return images


# ─────────────────────────────────────────────────────────────────
#  HTML → Editable DOCX (full pipeline)
# ─────────────────────────────────────────────────────────────────

def html_to_docx(
    html_path: str,
    output_path: str,
    page_size: str = "A4",
    orientation: str = "Portrait",
    margin_inches: float = 0.5,
    on_status: Optional[Callable[[str], None]] = None,
    keep_pdf: bool = False,
) -> tuple[str, Optional[str]]:
    """
    Full pipeline: HTML → PDF → Editable DOCX.
    Returns (docx_path, pdf_path_or_None).
    When keep_pdf=True, the intermediate PDF is preserved for preview.
    """
    status = on_status or (lambda _: None)

    # Create temp PDF
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        pdf_path = tmp.name

    try:
        render_to_pdf(
            html_path, pdf_path, page_size, orientation, margin_inches, on_status
        )
        pdf_to_docx(pdf_path, output_path, on_status)
    except Exception:
        try:
            os.unlink(pdf_path)
        except OSError:
            pass
        raise

    if not keep_pdf:
        try:
            os.unlink(pdf_path)
        except OSError:
            pass
        pdf_path = None

    return output_path, pdf_path


# ─────────────────────────────────────────────────────────────────
#  HTML → Screenshots (for PPTX)
# ─────────────────────────────────────────────────────────────────

async def _render_html_pages(
    html_path: str,
    viewport_width: int = 1200,
    device_scale_factor: float = 2.0,
    on_status: Optional[Callable[[str], None]] = None,
) -> list[bytes]:
    """
    Render HTML and screenshot each page container.
    Returns list of PNG byte buffers.
    """
    from playwright.async_api import async_playwright

    status = on_status or (lambda _: None)

    status("Launching browser...")
    async with async_playwright() as pw:
        browser = await _launch_best_browser(pw, headless=True)
        page = await browser.new_page(
            viewport={"width": viewport_width, "height": 1700},
            device_scale_factor=device_scale_factor,
        )

        abs_path = os.path.abspath(html_path)
        file_url = _build_file_url(abs_path)

        status("Loading HTML...")
        await page.goto(file_url, wait_until="networkidle")
        await page.evaluate("() => document.fonts.ready")

        # Remove floating UI elements
        await page.evaluate("""() => {
            const selectors = ['.fab-container', '.no-print', '[data-no-export]'];
            selectors.forEach(sel => {
                document.querySelectorAll(sel).forEach(el => el.remove());
            });
        }""")

        # Find page containers
        page_selectors = [".a4-page", ".page", "section", "article"]
        elements = []
        used_selector = None

        for sel in page_selectors:
            elements = await page.query_selector_all(sel)
            if elements:
                used_selector = sel
                break

        screenshots = []

        if elements:
            status(f"Found {len(elements)} pages (via '{used_selector}')...")
            for i, el in enumerate(elements):
                status(f"Capturing page {i + 1}/{len(elements)}...")
                png = await el.screenshot(type="png")
                screenshots.append(png)
        else:
            status("No page containers found — capturing full page...")
            png = await page.screenshot(type="png", full_page=True)
            screenshots.append(png)

        await browser.close()

    status(f"Captured {len(screenshots)} page(s).")
    return screenshots


def render_html_screenshots(
    html_path: str,
    dpi: int = 200,
    on_status: Optional[Callable[[str], None]] = None,
) -> list[bytes]:
    """Synchronous wrapper for screenshot rendering."""
    scale = dpi / 96.0
    return asyncio.run(
        _render_html_pages(
            html_path,
            viewport_width=1200,
            device_scale_factor=scale,
            on_status=on_status,
        )
    )


# ─────────────────────────────────────────────────────────────────
#  Screenshots → PPTX
# ─────────────────────────────────────────────────────────────────

def to_pptx(
    screenshots: list[bytes],
    output_path: str,
    slide_size: str = "Widescreen (16:9)",
    on_status: Optional[Callable[[str], None]] = None,
) -> str:
    """Build a PPTX with one image per slide."""
    from pptx import Presentation
    from pptx.util import Inches, Emu

    status = on_status or (lambda _: None)
    status("Building PPTX...")

    prs = Presentation()

    if "4:3" in slide_size:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    margin = Inches(0.3)
    content_w = slide_w - 2 * margin
    content_h = slide_h - 2 * margin

    blank_layout = prs.slide_layouts[6]

    for i, png in enumerate(screenshots):
        status(f"Adding slide {i + 1}/{len(screenshots)}...")

        slide = prs.slides.add_slide(blank_layout)

        img = Image.open(BytesIO(png))
        iw, ih = img.size
        aspect = ih / iw

        display_w = content_w
        display_h = int(display_w * aspect)

        if display_h > content_h:
            display_h = content_h
            display_w = int(display_h / aspect)

        left = (slide_w - display_w) // 2
        top = (slide_h - display_h) // 2

        slide.shapes.add_picture(
            BytesIO(png),
            left=left,
            top=top,
            width=display_w,
            height=display_h,
        )

    prs.save(output_path)
    status(f"PPTX saved: {os.path.basename(output_path)}")
    return output_path


def html_to_pptx(
    html_path: str,
    output_path: str,
    slide_size: str = "Widescreen (16:9)",
    dpi: int = 200,
    on_status: Optional[Callable[[str], None]] = None,
) -> str:
    """Full pipeline: HTML → Screenshots → PPTX."""
    screenshots = render_html_screenshots(html_path, dpi=dpi, on_status=on_status)
    if not screenshots:
        raise RuntimeError("No pages captured from HTML")
    return to_pptx(screenshots, output_path, slide_size=slide_size, on_status=on_status)


def html_to_editable_pptx(
    html_path: str,
    output_path: str,
    slide_size: str = "Widescreen (16:9)",
    on_status: Optional[Callable[[str], None]] = None,
) -> str:
    """Parse HTML to native editable PPTX elements."""
    import re
    from html.parser import HTMLParser
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    try:
        from pptx.enum.shapes import MSO_SHAPE
        RECTANGLE = MSO_SHAPE.RECTANGLE
    except ImportError:
        RECTANGLE = 1

    status = on_status or (lambda _: None)
    status("Parsing HTML to editable PPTX...")

    prs = Presentation()
    
    if "4:3" in slide_size:
        slide_w = Inches(10)
        slide_h = Inches(7.5)
    else:
        slide_w = Inches(13.333)
        slide_h = Inches(7.5)
        
    prs.slide_width = slide_w
    prs.slide_height = slide_h

    class PPTXHTMLParser(HTMLParser):
        def __init__(self):
            super().__init__()
            self.current_slide = None
            self.current_top = Inches(0.5)
            self.left = Inches(0.5)
            self.width = slide_w - Inches(1.0)
            
            self.tag_stack = []
            self.current_text = ""
            
            self.in_table = False
            self.table_data = []
            self.current_row = []
            self.is_header = False

        def _parse_style(self, style_str):
            styles = {}
            if not style_str: return styles
            for rule in style_str.split(';'):
                if ':' in rule:
                    k, v = rule.split(':', 1)
                    styles[k.strip().lower()] = v.strip()
            return styles

        def _parse_color(self, color_str):
            if not color_str: return None
            color_str = color_str.strip().lower()
            if color_str.startswith('#'):
                color_str = color_str.lstrip('#')
                if len(color_str) == 3:
                    color_str = "".join([c*2 for c in color_str])
                if len(color_str) in (6, 8):
                    try:
                        return RGBColor(int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16))
                    except ValueError:
                        pass
            elif color_str.startswith('rgb'):
                import re
                nums = re.findall(r'\d+', color_str)
                if len(nums) >= 3:
                    return RGBColor(int(nums[0]), int(nums[1]), int(nums[2]))
            return None

        def handle_starttag(self, tag, attrs):
            attr_dict = dict(attrs)
            classes = attr_dict.get("class", "")
            style = attr_dict.get("style", "")
            self.tag_stack.append((tag, classes, style))
            
            if tag == "div" and "a4-page" in classes:
                blank_layout = prs.slide_layouts[6]
                self.current_slide = prs.slides.add_slide(blank_layout)
                self.current_top = Inches(0.5)
                
                if "cover-page" in classes:
                    bg = self.current_slide.shapes.add_shape(
                        RECTANGLE, 0, 0, slide_w, slide_h
                    )
                    bg.fill.solid()
                    bg.fill.fore_color.rgb = RGBColor(100, 116, 139) # Slate 500
                    bg.line.fill.background()
                
            elif tag == "table":
                self.in_table = True
                self.table_data = []
                
            elif tag == "tr" and self.in_table:
                self.current_row = []
                
            elif tag in ("td", "th") and self.in_table:
                self.current_text = ""
                self.is_header = (tag == "th")

        def handle_endtag(self, tag):
            if not self.tag_stack: return
            
            popped_tag = None
            classes = ""
            style = ""
            # Pop up to the matching tag (scan right-to-left for the most
            # recent open of this tag; pop() already removes it in-place)
            for i in range(len(self.tag_stack) - 1, -1, -1):
                if self.tag_stack[i][0] == tag:
                    popped_tag, classes, style = self.tag_stack.pop(i)
                    break
                    
            if not popped_tag: return
            
            text = self.current_text.strip()
            self.current_text = ""
            
            if not self.current_slide:
                return
                
            if tag in ("td", "th") and self.in_table:
                self.current_row.append((text, self.is_header))
                
            elif tag == "tr" and self.in_table:
                self.table_data.append(self.current_row)
                
            elif tag == "table":
                self.in_table = False
                self._render_table()
                
            elif tag in ("h1", "h2", "h3", "h4", "p", "li", "div") and not self.in_table:
                if text:
                    self._render_text(tag, text, classes, style)

        def handle_data(self, data):
            if self.current_slide:
                # Replace multiple spaces/newlines with a single space
                clean_data = re.sub(r'\s+', ' ', data)
                self.current_text += clean_data

        def _render_text(self, tag, text, classes, style_str=""):
            font_size = Pt(12)
            bold = False
            italic = False
            color = RGBColor(51, 65, 85)
            bg_color = None
            border_color = None
            height = Inches(0.4)
            
            if tag == "h1":
                font_size = Pt(36)
                bold = True
                color = RGBColor(15, 23, 42)
                height = Inches(0.8)
            elif tag == "h2":
                font_size = Pt(28)
                bold = True
                color = RGBColor(15, 23, 42)
                height = Inches(0.6)
            elif tag == "h3":
                font_size = Pt(20)
                bold = True
                color = RGBColor(71, 85, 105)
                height = Inches(0.5)
            elif tag == "h4":
                font_size = Pt(16)
                bold = True
                height = Inches(0.4)
            elif tag == "li":
                text = "• " + text
                
            styles = self._parse_style(style_str)
            if 'font-weight' in styles:
                val = styles['font-weight']
                if val in ('bold', 'bolder') or (val.isdigit() and int(val) >= 600):
                    bold = True
                elif val in ('normal', 'lighter') or (val.isdigit() and int(val) < 600):
                    bold = False
            if 'font-style' in styles and styles['font-style'] == 'italic':
                italic = True
            if 'color' in styles:
                parsed_color = self._parse_color(styles['color'])
                if parsed_color: color = parsed_color
            if 'background-color' in styles:
                bg_color = self._parse_color(styles['background-color'])
            if 'border' in styles or 'border-color' in styles:
                border_color = self._parse_color(styles.get('border-color') or styles.get('border'))
                
            if any("cover-page" in c for t, c, s in self.tag_stack) or "cover-page" in classes:
                color = RGBColor(255, 255, 255)
                
            # Prevent shapes from going completely off slide
            if self.current_top > slide_h - Inches(0.5):
                return
                
            if bg_color or border_color:
                txBox = self.current_slide.shapes.add_shape(RECTANGLE, self.left, self.current_top, self.width, height)
                if bg_color:
                    txBox.fill.solid()
                    txBox.fill.fore_color.rgb = bg_color
                else:
                    txBox.fill.background()
                    
                if border_color:
                    txBox.line.color.rgb = border_color
                    txBox.line.width = Pt(1)
                else:
                    txBox.line.fill.background()
            else:
                txBox = self.current_slide.shapes.add_textbox(self.left, self.current_top, self.width, height)
                
            tf = txBox.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = font_size
            p.font.bold = bold
            p.font.italic = italic
            p.font.color.rgb = color
            
            chars_per_line = 90 if font_size.pt < 18 else 45
            lines = max(1, math.ceil(len(text) / chars_per_line))
            self.current_top += Inches((font_size.pt / 72.0) * lines * 1.5) + Inches(0.1)

        def _render_table(self):
            if not self.table_data: return
            rows = len(self.table_data)
            cols = max(len(r) for r in self.table_data) if rows > 0 else 0
            if cols == 0: return
            
            height = Inches(0.4 * rows)
            if self.current_top + height > slide_h:
                height = max(Inches(0.5), slide_h - self.current_top - Inches(0.2))
            
            try:
                table_shape = self.current_slide.shapes.add_table(
                    rows, cols, self.left, self.current_top, self.width, height
                ).table
                
                for r_idx, row in enumerate(self.table_data):
                    for c_idx, (text, is_header) in enumerate(row):
                        if c_idx < cols:
                            cell = table_shape.cell(r_idx, c_idx)
                            cell.text = text
                            p = cell.text_frame.paragraphs[0]
                            p.font.size = Pt(10)
                            if is_header:
                                p.font.bold = True
                self.current_top += height + Inches(0.2)
            except Exception as e:
                print(f"Table error ignored: {e}")

    with open(html_path, 'r', encoding='utf-8', errors='ignore') as f:
        html_content = f.read()
        
    parser = PPTXHTMLParser()
    parser.feed(html_content)
    
    if len(prs.slides) == 0:
        # Fallback if no pages found
        blank_layout = prs.slide_layouts[6]
        prs.slides.add_slide(blank_layout)
        
    prs.save(output_path)
    status(f"Editable PPTX saved: {os.path.basename(output_path)}")
    return output_path


# ─────────────────────────────────────────────────────────────────
#  Playwright setup helpers
# ─────────────────────────────────────────────────────────────────

def check_playwright_installed() -> bool:
    """
    Check if any usable browser is available.
    Tries system Chrome, Edge, then Playwright's own Chromium.
    """
    try:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as pw:
            for channel in _BROWSER_CHANNELS:
                try:
                    opts = {"headless": True}
                    if channel:
                        opts["channel"] = channel
                    browser = pw.chromium.launch(**opts)
                    browser.close()
                    return True
                except Exception:
                    continue
        return False
    except Exception:
        return False


def get_available_browser_name() -> str:
    """
    Return the name of the first usable browser, e.g. 'Chrome', 'Edge',
    or 'Playwright Chromium'.  Returns '' if none found.
    """
    labels = {"chrome": "Chrome", "msedge": "Edge", None: "Playwright Chromium"}
    try:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as pw:
            for channel in _BROWSER_CHANNELS:
                try:
                    opts = {"headless": True}
                    if channel:
                        opts["channel"] = channel
                    browser = pw.chromium.launch(**opts)
                    browser.close()
                    return labels[channel]
                except Exception:
                    continue
    except Exception:
        pass
    return ""


def install_playwright_browser(
    on_status: Optional[Callable[[str], None]] = None,
) -> bool:
    """
    Last-resort: install Playwright's bundled Chromium.
    Only reached when no system browser (Chrome/Edge) is found.
    """
    import subprocess

    status = on_status or (lambda _: None)
    status("No system browser found. Installing Chromium (~150MB)...")

    # When frozen, sys.executable IS the exe — cannot use it as python.
    # Fall back to the playwright CLI that ships with the Python package.
    if getattr(sys, "frozen", False):
        status("Run: python -m playwright install chromium")
        status("Or install Google Chrome or Microsoft Edge.")
        return False

    try:
        result = subprocess.run(
            [sys.executable, "-m", "playwright", "install", "chromium"],
            capture_output=True,
            text=True,
            timeout=300,
        )
        if result.returncode == 0:
            status("Chromium installed successfully.")
            return True
        else:
            status(f"Install failed: {result.stderr}")
            return False
    except Exception as e:
        status(f"Install error: {e}")
        return False


def pdf_to_editable_pptx(
    pdf_path: str,
    output_path: str,
    slide_size: str = "Widescreen (16:9)",
    on_status: Optional[Callable[[str], None]] = None,
) -> str:
    """Convert PDF pages to maximally editable PPTX objects using scene graph."""
    from pdf_extract import extract_pdf_to_scene_graph
    from pptx_writer import write_scene_graph_to_pptx
    
    # Phase 1: PDF Extraction -> Scene Graph
    pages, warnings = extract_pdf_to_scene_graph(pdf_path, on_status=on_status)
    
    # Phase 2: Scene Graph -> PPTX Writer
    return write_scene_graph_to_pptx(pages, output_path, slide_size=slide_size, warnings=warnings, on_status=on_status)
